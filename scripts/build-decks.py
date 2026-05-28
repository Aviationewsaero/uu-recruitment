#!/usr/bin/env python3
"""
Generate two PPTX decks reflecting the live state of careers.ews.aero:

  1. UU-Aviation-Drive-Day-Playbook.pptx   — for the on-site EWS team
  2. UU-Aviation-Super-Admin-Guide.pptx    — for the founder / super admin

Run:
    python3 scripts/build-decks.py
Outputs land in /docs/.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── Brand ────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x8A)
NAVY_DARK   = RGBColor(0x17, 0x2A, 0x5E)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)
GREEN_DARK  = RGBColor(0x16, 0xA3, 0x4A)
BG          = RGBColor(0xF8, 0xFA, 0xFC)
SURFACE     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)
MUTED       = RGBColor(0x64, 0x74, 0x8B)
TEXT        = RGBColor(0x0F, 0x17, 0x2A)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
RED         = RGBColor(0xDC, 0x26, 0x26)
BLUE        = RGBColor(0x25, 0x63, 0xEB)

OUT_DIR = Path(__file__).resolve().parent.parent / "docs"
OUT_DIR.mkdir(parents=True, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # solid white background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = SURFACE
    return slide


def add_rect(slide, x, y, w, h, fill, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.line.fill.background() if line is None else None
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is not None:
        r.line.color.rgb = line
        r.line.width = Pt(0.75)
    return r


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT,
                line_spacing=1.15, bullet_color=GREEN_DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2);  tf.margin_bottom = Pt(2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
        # bullet glyph
        rb = p.add_run()
        rb.text = "■  "
        rb.font.name = "Calibri"
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rb.font.bold = True
        # body text — supports:
        #   "plain string"                  → single run
        #   ("Label", "body text")          → bold label, normal body separated by " — "
        #   [("chunk", True), (".", False)] → fine-grained run list
        if isinstance(item, tuple) and len(item) == 2 and all(isinstance(x, str) for x in item):
            label, body = item
            rl = p.add_run(); rl.text = label
            rl.font.name = "Calibri"; rl.font.size = Pt(size)
            rl.font.bold = True; rl.font.color.rgb = color
            rsep = p.add_run(); rsep.text = " — "
            rsep.font.name = "Calibri"; rsep.font.size = Pt(size)
            rsep.font.color.rgb = MUTED
            rb = p.add_run(); rb.text = body
            rb.font.name = "Calibri"; rb.font.size = Pt(size)
            rb.font.color.rgb = color
        elif isinstance(item, list):
            for chunk, is_bold in item:
                rr = p.add_run()
                rr.text = chunk
                rr.font.name = "Calibri"
                rr.font.size = Pt(size)
                rr.font.bold = is_bold
                rr.font.color.rgb = color
        else:
            rr = p.add_run()
            rr.text = str(item)
            rr.font.name = "Calibri"
            rr.font.size = Pt(size)
            rr.font.color.rgb = color
    return tb


def page_header(slide, kicker, title, subtitle=None):
    """Compact navy header bar at the top of every content slide."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), NAVY)
    add_text(slide, Inches(0.55), Inches(0.18), Inches(8), Inches(0.3),
             kicker, size=10, bold=True, color=GREEN,
             font="Calibri")
    add_text(slide, Inches(0.55), Inches(0.40), Inches(11), Inches(0.55),
             title, size=22, bold=True, color=SURFACE, font="Georgia")
    if subtitle:
        add_text(slide, Inches(0.55), Inches(0.78), Inches(12), Inches(0.3),
                 subtitle, size=11, color=RGBColor(0xCB, 0xD5, 0xE1))
    # subtle accent line
    add_rect(slide, 0, Inches(1.0), SLIDE_W, Pt(2), GREEN)


def footer(slide, page_num, total, doc_title):
    add_text(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3),
             doc_title, size=9, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
             f"{page_num} / {total}", size=9, color=MUTED,
             align=PP_ALIGN.RIGHT)


# ─── Title-slide layout ──────────────────────────────────────────────────

def title_slide(prs, kicker, title, subtitle, audience):
    s = blank_slide(prs)
    # full-bleed navy
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # diagonal green accent
    accent = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                Inches(10.5), Inches(4.5),
                                Inches(2.8), Inches(3.0))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.rotation = 180

    add_text(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.4),
             "ELITE WORLD SERVICES · AVIATION RECRUITMENT",
             size=11, bold=True, color=GREEN, font="Calibri")
    add_text(s, Inches(0.7), Inches(1.2), Inches(11), Inches(0.6),
             kicker, size=14, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1),
             font="Calibri")
    add_text(s, Inches(0.7), Inches(2.0), Inches(11), Inches(2.0),
             title, size=48, bold=True, color=SURFACE, font="Georgia")
    add_text(s, Inches(0.7), Inches(4.2), Inches(11), Inches(1.2),
             subtitle, size=16, color=RGBColor(0xCB, 0xD5, 0xE1),
             font="Calibri")
    # audience pill
    pill = add_rect(s, Inches(0.7), Inches(5.8), Inches(5.5), Inches(0.6),
                    GREEN)
    add_text(s, Inches(0.85), Inches(5.86), Inches(5.2), Inches(0.5),
             f"FOR: {audience.upper()}", size=12, bold=True,
             color=NAVY, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)

    # bottom meta
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             "careers.ews.aero  ·  Live system snapshot  ·  Updated 2026-05-28",
             size=10, color=RGBColor(0xCB, 0xD5, 0xE1), font="Calibri")
    return s


# ─── Section divider ─────────────────────────────────────────────────────

def section_slide(prs, number, title, hint):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    add_rect(s, 0, Inches(3.2), SLIDE_W, Pt(2), GREEN)
    add_text(s, Inches(0.7), Inches(2.0), Inches(2.2), Inches(1.1),
             f"PART {number:02d}", size=14, bold=True, color=GREEN,
             font="Calibri")
    add_text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(1.0),
             title, size=42, bold=True, color=SURFACE, font="Georgia")
    add_text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(0.5),
             hint, size=14, color=RGBColor(0xCB, 0xD5, 0xE1), font="Calibri")
    return s


# ─── 3-card row (used for step-by-step + role cards) ──────────────────────

def three_card_slide(prs, kicker, title, cards, *, footer_text=None):
    """
    cards = [{title, body (str or [bullets]), accent: RGBColor}]
    """
    s = blank_slide(prs)
    page_header(s, kicker, title)
    y = Inches(1.5)
    gap = Inches(0.35)
    card_w = (SLIDE_W - Inches(1.1) - gap * (len(cards) - 1)) / len(cards)
    x = Inches(0.55)
    for c in cards:
        accent = c.get("accent", GREEN)
        # accent bar
        add_rect(s, x, y, card_w, Inches(0.18), accent)
        # body card
        body = add_rect(s, x, y + Inches(0.18), card_w, Inches(5.0), SURFACE,
                        line=BORDER)
        # title
        add_text(s, x + Inches(0.25), y + Inches(0.35),
                 card_w - Inches(0.5), Inches(0.6),
                 c["title"], size=16, bold=True, color=NAVY, font="Georgia")
        # body
        body_y = y + Inches(1.0)
        if isinstance(c["body"], list):
            add_bullets(s, x + Inches(0.25), body_y,
                        card_w - Inches(0.5), Inches(4.2),
                        c["body"], size=12, bullet_color=accent)
        else:
            add_text(s, x + Inches(0.25), body_y,
                     card_w - Inches(0.5), Inches(4.2),
                     c["body"], size=12, color=TEXT)
        x += card_w + gap
    if footer_text:
        add_text(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.35),
                 footer_text, size=11, color=MUTED, font="Calibri")
    return s


# ─── Two-column slide ────────────────────────────────────────────────────

def two_col_slide(prs, kicker, title, left_title, left_items,
                  right_title, right_items, *, right_accent=GREEN,
                  left_accent=BLUE, footer_text=None):
    s = blank_slide(prs)
    page_header(s, kicker, title)
    gap = Inches(0.35)
    col_w = (SLIDE_W - Inches(1.1) - gap) / 2

    for x, h_title, items, accent in [
        (Inches(0.55), left_title, left_items, left_accent),
        (Inches(0.55) + col_w + gap, right_title, right_items, right_accent),
    ]:
        add_rect(s, x, Inches(1.5), col_w, Inches(0.18), accent)
        add_rect(s, x, Inches(1.68), col_w, Inches(5.0), SURFACE, line=BORDER)
        add_text(s, x + Inches(0.25), Inches(1.85), col_w - Inches(0.5),
                 Inches(0.5), h_title, size=16, bold=True, color=NAVY,
                 font="Georgia")
        add_bullets(s, x + Inches(0.25), Inches(2.45), col_w - Inches(0.5),
                    Inches(4.0), items, size=12, bullet_color=accent)
    if footer_text:
        add_text(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.35),
                 footer_text, size=11, color=MUTED, font="Calibri")
    return s


# ─── Step-list slide (vertical numbered) ─────────────────────────────────

def steps_slide(prs, kicker, title, steps, *, footer_text=None):
    """steps = [(label, body)]"""
    s = blank_slide(prs)
    page_header(s, kicker, title)
    y = Inches(1.4)
    row_h = (Inches(5.2)) / len(steps)
    for i, (label, body) in enumerate(steps):
        ry = y + row_h * i
        # number bubble
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55),
                                  ry + Inches(0.05), Inches(0.55),
                                  Inches(0.55))
        circ.line.fill.background()
        circ.fill.solid()
        circ.fill.fore_color.rgb = NAVY
        add_text(s, Inches(0.55), ry + Inches(0.05),
                 Inches(0.55), Inches(0.55), str(i + 1),
                 size=18, bold=True, color=SURFACE, font="Georgia",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # label
        add_text(s, Inches(1.3), ry + Inches(0.0), Inches(11), Inches(0.4),
                 label, size=15, bold=True, color=NAVY, font="Georgia")
        # body
        add_text(s, Inches(1.3), ry + Inches(0.4), Inches(11.5),
                 row_h - Inches(0.4), body, size=12, color=TEXT)
    if footer_text:
        add_text(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.35),
                 footer_text, size=11, color=MUTED, font="Calibri")
    return s


# ─── Table slide (issue ↔ fix) ───────────────────────────────────────────

def issue_fix_slide(prs, kicker, title, rows, *, col_titles=("Symptom", "Cause", "Do this"),
                    footer_text=None):
    s = blank_slide(prs)
    page_header(s, kicker, title)
    x = Inches(0.55)
    y = Inches(1.45)
    w = SLIDE_W - Inches(1.1)
    cols = [0.30, 0.27, 0.43]
    col_w = [w * c for c in cols]
    # header
    add_rect(s, x, y, w, Inches(0.45), NAVY)
    cx = x
    for i, t in enumerate(col_titles):
        add_text(s, cx + Inches(0.18), y + Inches(0.08),
                 col_w[i] - Inches(0.2), Inches(0.4),
                 t.upper(), size=10, bold=True, color=SURFACE,
                 font="Calibri")
        cx += col_w[i]
    # body rows
    ry = y + Inches(0.45)
    row_h = Inches(0.90)
    for i, (sym, cause, fix) in enumerate(rows):
        fill = SURFACE if i % 2 == 0 else BG
        add_rect(s, x, ry, w, row_h, fill, line=BORDER)
        cx = x
        for j, txt in enumerate((sym, cause, fix)):
            color = NAVY if j == 0 else (MUTED if j == 1 else TEXT)
            bold = j == 0
            add_text(s, cx + Inches(0.18), ry + Inches(0.10),
                     col_w[j] - Inches(0.2), row_h - Inches(0.2),
                     txt, size=11, bold=bold, color=color)
            cx += col_w[j]
        ry += row_h
    if footer_text:
        add_text(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.35),
                 footer_text, size=10.5, color=MUTED, font="Calibri")
    return s


# ─── Closing slide ───────────────────────────────────────────────────────

def closing_slide(prs, title, lines, contacts=None):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(2.7), SLIDE_W, Pt(2), GREEN)
    add_text(s, Inches(0.7), Inches(1.4), Inches(11), Inches(1.2),
             title, size=42, bold=True, color=SURFACE, font="Georgia")
    add_text(s, Inches(0.7), Inches(3.1), Inches(12), Inches(0.5),
             lines[0] if lines else "",
             size=18, color=GREEN, font="Calibri", bold=True)
    if len(lines) > 1:
        add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(2.0),
                 "\n".join(lines[1:]),
                 size=13, color=RGBColor(0xCB, 0xD5, 0xE1),
                 font="Calibri")
    if contacts:
        add_text(s, Inches(0.7), Inches(6.0), Inches(2), Inches(0.4),
                 "QUICK CONTACTS", size=10, bold=True, color=GREEN,
                 font="Calibri")
        add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(1.0),
                 "    ·    ".join(contacts), size=11,
                 color=RGBColor(0xCB, 0xD5, 0xE1), font="Calibri")
    return s


# ═══════════════════════════════════════════════════════════════════════════
#  DECK 1 — DRIVE-DAY PLAYBOOK (for the on-site EWS team)
# ═══════════════════════════════════════════════════════════════════════════

def build_playbook():
    prs = new_prs()
    DOC_TITLE = "UU Aviation Recruitment 2026  ·  Drive-Day Playbook"
    slides = []

    # ── Cover ────────────────────────────────────────────────────────────
    slides.append(title_slide(
        prs,
        kicker="DRIVE-DAY PLAYBOOK  ·  EDITION 2026.05",
        title="The drive runs itself.\nYour job is to keep it calm.",
        subtitle=(
            "Step-by-step guide for desk operators, recruiters and floor team. "
            "Covers the entire student journey — arrival to result — and every "
            "issue we've seen in testing."
        ),
        audience="On-site EWS team (desk operators · recruiters · floor)",
    ))

    # ── TOC / what's inside ──────────────────────────────────────────────
    s = blank_slide(prs)
    page_header(s, "INSIDE THIS PLAYBOOK", "Eight parts. Read in order on day-1.")
    parts = [
        ("01", "Before the doors open", "T-60 min · venue & system pre-flight"),
        ("02", "How a student moves through the venue", "End-to-end flow on one page"),
        ("03", "Roles at the desk", "Who does what · how to swap places"),
        ("04", "Calling tokens & running interviews", "The recruiter loop"),
        ("05", "Crowd control & no-shows", "The desk operator playbook"),
        ("06", "Live monitoring", "The super admin's dashboard during the drive"),
        ("07", "When things go wrong", "Symptom → cause → fix"),
        ("08", "Wrap-up & exports", "Selection list, CSV, photo backup"),
    ]
    y = Inches(1.45)
    for i, (num, title, hint) in enumerate(parts):
        row_y = y + Inches(0.62) * i
        add_rect(s, Inches(0.55), row_y, Inches(0.55), Inches(0.50), NAVY)
        add_text(s, Inches(0.55), row_y, Inches(0.55), Inches(0.50),
                 num, size=14, bold=True, color=SURFACE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Georgia")
        add_text(s, Inches(1.3), row_y + Inches(0.02),
                 Inches(7.5), Inches(0.3), title,
                 size=14, bold=True, color=NAVY, font="Georgia")
        add_text(s, Inches(1.3), row_y + Inches(0.28),
                 Inches(11), Inches(0.3), hint,
                 size=11, color=MUTED)
    slides.append(s)

    # ─── PART 01 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 1, "Before the doors open",
        "T-60 minutes. 8 checks. If all green, you're ready."))

    slides.append(steps_slide(
        prs, "PART 01  ·  PRE-FLIGHT", "8-step venue + system check",
        [
            ("1.  Power up the TV / display monitor",
             "Open Chrome → https://careers.ews.aero/display  →  press F11 for full-screen. "
             "Token board should show the room cards. Auto-refreshes every 5 sec."),
            ("2.  Hang the QR poster at the entrance",
             "Print from /admin/qr-poster (super admin only). One QR per entrance is enough. "
             "Students scan → land directly on /register."),
            ("3.  Recruiters log in at their laptops",
             "https://careers.ews.aero/admin/login  →  email + password (no OTP). "
             "Each recruiter sees their assigned room auto-selected."),
            ("4.  Desk operator logs in",
             "Same URL, same flow. Desk operator's home screen is the live queue at /admin/queue."),
            ("5.  Open the runbook for quick reference",
             "/admin/runbook  →  print 1 copy and keep at the desk. Same content as this deck, "
             "compressed to a single A4 sheet."),
            ("6.  Sanity-check the live tiles",
             "Open /admin in another tab. Registered/Waiting/Done numbers should match reality. "
             "If they don't refresh — check Vercel status (super admin)."),
            ("7.  Brief the floor team (volunteers)",
             "Two volunteer roles: 'Greeter' (helps students scan QR), "
             "'Caller' (walks the waiting area calling token numbers from the TV)."),
            ("8.  Open this playbook bookmark on every device",
             "https://careers.ews.aero/admin/runbook  — the runbook page mirrors this deck."),
        ],
        footer_text="If any of these 8 steps fail → call the super admin BEFORE opening doors. "
                    "Do not improvise."))

    # ─── PART 02 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 2, "How a student moves through the venue",
        "Six stops. Roughly 8–12 minutes per student end-to-end."))

    slides.append(steps_slide(
        prs, "PART 02  ·  THE STUDENT JOURNEY",
        "From entrance to result — one screen at a time",
        [
            ("Stop 1  ·  Entrance",
             "Greeter points to the QR poster. Student scans → phone opens "
             "https://careers.ews.aero/register."),
            ("Stop 2  ·  Registration form (5 min)",
             "Student enters details + uploads resume + photo (compressed automatically "
             "in-browser so phone-camera photos don't fail). Submit → token issued."),
            ("Stop 3  ·  Token & admit card",
             "Success screen shows BIG green token number. Student downloads the PDF admit card "
             "(also emailed). They keep this on their phone."),
            ("Stop 4  ·  Waiting area",
             "Student watches the TV. Display shows which room is calling which token. "
             "Updates every 5 sec."),
            ("Stop 5  ·  Interview",
             "When their number shows up next to a room, they walk to that room. "
             "Recruiter has their full profile + resume + photo on screen."),
            ("Stop 6  ·  Result",
             "Recruiter saves decision → student gets an automatic email with the outcome. "
             "They can leave the venue."),
        ],
        footer_text="No paper forms. No phone calls. No SMS. Everything is on the TV "
                    "or in email."))

    # ─── PART 03 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 3, "Roles at the desk",
        "Four roles. Each one has a screen. None of them step on each other."))

    slides.append(three_card_slide(
        prs, "PART 03  ·  ROLES — TIER 1", "Recruiter & Desk Operator (the busy ones)",
        [
            {"title": "RECRUITER", "accent": GREEN, "body": [
                "Logs in at /admin/login. Sees ONLY their room.",
                "Clicks 'Call Next Token' to summon the next student.",
                "Reviews the student profile on screen, conducts interview.",
                "Records: rating 1–5, notes, decision (shortlist / hold / etc.).",
                "Toggles 'Auto-advance' → call the next one immediately after save.",
                "Cannot see or modify other rooms.",
            ]},
            {"title": "DESK OPERATOR", "accent": BLUE, "body": [
                "Logs in at /admin/login. Lands at the LIVE QUEUE.",
                "Watches the full waiting list across all rooms.",
                "Helps no-shows: re-issues skipped tokens via 'Recall'.",
                "Marks permanent no-shows so they stop appearing in the queue.",
                "Cannot record interview decisions — only manages flow.",
                "First point of contact for confused or upset students.",
            ]},
            {"title": "FLOOR (no login)", "accent": AMBER, "body": [
                "Greeter at entrance: helps students scan QR, hands out a backup paper QR if their phone is dead.",
                "Caller in waiting area: reads token numbers off the TV every 30 sec for students who missed it.",
                "Runner: physically walks no-shows from the waiting area to the recruiter room.",
            ]},
        ]))

    slides.append(three_card_slide(
        prs, "PART 03  ·  ROLES — TIER 2", "Email Manager & Super Admin (the office roles)",
        [
            {"title": "EMAIL MANAGER", "accent": GREEN, "body": [
                "Sends post-drive announcements to filtered groups (e.g. all SELECTED, or all REJECTED).",
                "Uses /admin/emails: pick audience → preview count → type confirmation if > 25 → send.",
                "Does not touch the queue, rooms, or interviews.",
                "Can export the full student list as CSV from /admin/students.",
            ]},
            {"title": "SUPER ADMIN (founder)", "accent": NAVY, "body": [
                "Creates all staff accounts at /admin/users.",
                "Assigns each RECRUITER to a room from the same page (dropdown).",
                "Watches live analytics + tile counts at /admin.",
                "Has access to ALL other panels for support / firefighting.",
                "Only role that can use the emergency env-var password.",
            ]},
            {"title": "ROLE SWAPS", "accent": MUTED, "body": [
                "All staff log in with the SAME URL — only the role differs.",
                "If a recruiter calls in sick: super admin reassigns their room to a stand-by recruiter from /admin/users.",
                "Desk operator backup: any super admin can use /admin/queue.",
                "Never share a login. Make a new account in 30 seconds.",
            ]},
        ]))

    # ─── PART 04 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 4, "Calling tokens & running interviews",
        "The recruiter screen — what every click does."))

    slides.append(steps_slide(
        prs, "PART 04  ·  THE RECRUITER LOOP", "6 clicks per student",
        [
            ("1.  Click 'Call Next Token'",
             "Pulls the lowest WAITING token into your room. The student sees their token appear "
             "next to your room name on the TV."),
            ("2.  Wait ~30 sec for the student to walk in",
             "If they don't show up: click 'Skip (no-show)'. Token goes to SKIPPED. Desk operator can recall later."),
            ("3.  Click 'Start interview'",
             "Marks the token IN_PROGRESS. Begins timing the interview (used in analytics)."),
            ("4.  Conduct the interview",
             "Student profile, resume preview, and photo are right there. Take notes in the box — "
             "they're saved with the decision."),
            ("5.  Pick a decision + rating",
             "5 options: Shortlist · Hold · Re-interview · Reject · SELECTED. Rating is 1-5 stars. "
             "Notes are optional but encouraged."),
            ("6.  Click 'Save & Next'",
             "Decision is locked. Student gets an automatic email with their result. "
             "If 'Auto-advance' is on, your screen instantly calls the next token."),
        ],
        footer_text="Average interview: 8 min. With auto-advance on, you handle ~7 students/hour."))

    # ─── PART 05 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 5, "Crowd control & no-shows",
        "The desk operator's playbook — keep the line moving."))

    slides.append(two_col_slide(
        prs, "PART 05  ·  DESK OPERATOR", "Live-queue mastery",
        left_title="WHEN TO INTERVENE",
        left_items=[
            "Token waiting > 20 min → check if the student is still in the venue",
            "A room is idle but the queue is long → call that recruiter on the phone",
            "Same name registered twice → keep the latest token, no-show the older one",
            "Student says they didn't get the email → resend from /admin/students/[regId]",
            "Student says 'I can't find my token number' → search by name in /admin/students",
        ],
        right_title="THE 4 BUTTONS YOU USE",
        right_items=[
            ("Recall", "Brings a SKIPPED token back to the front of WAITING. One click."),
            ("No-show", "Permanent — token disappears from the queue. Use only after 2 callbacks."),
            ("Filter chip 'Skipped'", "Shows just the skipped tokens. Easiest to recall in batch."),
            ("Search by token #", "Top of the queue page. Type the number → jump to row."),
        ],
        footer_text="Rule: never delete a token. Use 'No-show'. Deletion would break analytics + the student's email trail."))

    # ─── PART 06 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 6, "Live monitoring",
        "The super admin's dashboard — what to watch, what to ignore."))

    slides.append(two_col_slide(
        prs, "PART 06  ·  WATCH THESE NUMBERS", "Real-time dashboard at /admin",
        left_title="GREEN — KEEP DOING IT",
        left_items=[
            "Registered count climbing steadily (~2/min during arrival rush)",
            "Waiting < In-progress × 4 → flow is balanced",
            "All 4 room cards on /display showing a 'Now serving' badge",
            "Completed count rising at ~1 every 90 sec (across all rooms)",
            "No 'Failed' emails in the audit log",
        ],
        right_title="AMBER / RED — INVESTIGATE",
        right_items=[
            ("AMBER", "Waiting > 50 students → ask if a recruiter is stuck"),
            ("AMBER", "One room hasn't called a token in 15 min → recruiter break or laptop issue"),
            ("RED", "Tiles stopped updating → page is stale, refresh the browser"),
            ("RED", "Registered count went BACKWARDS → call the dev immediately (data integrity issue)"),
            ("RED", "Display TV shows '—' for all rooms → DB or Vercel outage, check status pages"),
        ],
        footer_text="The /admin overview page now loads in ~350 ms warm. If it spins for > 5 sec → something is wrong (not just slow)."))

    # ─── PART 07 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 7, "When things go wrong",
        "Twelve scenarios we've seen in testing. Fix flow under each."))

    slides.append(issue_fix_slide(
        prs, "PART 07  ·  TROUBLESHOOTING — STUDENT SIDE",
        "Most common issues at the entrance and registration",
        [
            ("QR code won't scan",
             "Phone camera dirty or QR too small from distance",
             "Hand them the paper QR card (printed copies kept at desk) — same URL"),
            ("Form rejects their phone photo",
             "File over 2 MB after compression",
             "Tell them to retake with lower resolution OR upload from gallery. Compression is automatic."),
            ("'Email already registered' error",
             "They tried earlier — their token already exists",
             "Search by email in /admin/students → give them the existing token number"),
            ("OTP email never arrives",
             "OTP is OFF in current config — should never see this",
             "If they do see it: super admin pushes APP_STUDENT_AUTH=none again on Vercel"),
            ("Admit card PDF won't open on their phone",
             "Older Android browsers don't render Vercel-served PDFs inline",
             "They can still attend with just the token number. Email also has the number."),
            ("Submit form spins forever",
             "Their phone signal dropped mid-upload",
             "Re-submit when signal returns. Server checks for duplicate email — won't double-register."),
        ],
        footer_text="Rate-limit: a single phone is capped at 5 submissions / 10 min. Don't worry — this is to block bots, not a real student issue."))

    slides.append(issue_fix_slide(
        prs, "PART 07  ·  TROUBLESHOOTING — STAFF SIDE",
        "Issues recruiters and desk operators may hit",
        [
            ("Recruiter sees 'No room assigned to you yet'",
             "Their account exists but Room.recruiterId isn't set",
             "Super admin opens /admin/users → finds their row → picks a room from the dropdown → done. Recruiter refreshes."),
            ("Login says 'Too many attempts'",
             "10 failed logins from this IP in 5 min — 15-min cooldown kicked in",
             "Wait it out. Or: super admin uses a different device/network OR the emergency password."),
            ("Decision form rejects 'Token is not in this room'",
             "Their room reassignment happened mid-interview, OR a stale tab",
             "Refresh the recruiter page. The new room context loads correctly."),
            ("Display board frozen for everyone",
             "TV browser tab crashed or lost network",
             "Refresh the TV browser (Cmd/Ctrl+R). It auto-resumes."),
            ("Bulk-email send won't trigger",
             "Audience > 25 — typed confirmation required (anti-mistake guard)",
             "Type the EXACT number shown in the recipient count, then click Send."),
            ("CSV export downloads empty file",
             "Filter selected zero students — there's nothing to export",
             "Clear all filters → re-export. Headers always present even if zero rows."),
        ],
        footer_text="Every recoverable error tells you what to do in the error message itself. Read it before calling for help."))

    # ─── PART 08 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 8, "Wrap-up & exports",
        "After the last student walks out — what to deliver."))

    slides.append(steps_slide(
        prs, "PART 08  ·  END-OF-DAY", "5 things to do before you go home",
        [
            ("1.  Confirm zero WAITING tokens",
             "Open /admin/queue → filter 'Waiting'. Should be empty. "
             "If not: are they actually here? No-show or recall."),
            ("2.  Run the selection email blast",
             "Email Manager opens /admin/emails → audience: SELECTED → write congrats email → "
             "type recipient count to confirm → Send."),
            ("3.  Export the full student CSV",
             "/admin/students → 'Export CSV' top right. Includes ALL students, all statuses. "
             "Keep this file safe — it's your audit trail."),
            ("4.  Take a screenshot of /admin tiles",
             "The final numbers for the day. Goes in the post-drive report."),
            ("5.  Hand over to the super admin",
             "They run the post-drive backup script + archive the Supabase data."),
        ]))

    # ── Closing ─────────────────────────────────────────────────────────
    slides.append(closing_slide(
        prs,
        "You've got this.",
        [
            "The system handles 700 students. Your job is to keep humans calm.",
            "",
            "When in doubt: read the runbook printout at the desk.",
            "When the runbook doesn't help: call the super admin.",
            "When the super admin doesn't pick up: the system keeps running even if you click nothing — students keep registering, tokens keep flowing.",
        ],
        contacts=[
            "Super admin: bhupender@eliteworldservices.com",
            "Live URL: careers.ews.aero",
            "Runbook: careers.ews.aero/admin/runbook",
        ]))

    # ── Footers ──────────────────────────────────────────────────────────
    total = len(slides)
    for i, s in enumerate(slides[2:], start=2):
        footer(s, i + 1, total, DOC_TITLE)

    out = OUT_DIR / "UU-Aviation-Drive-Day-Playbook.pptx"
    prs.save(out)
    return out


# ═══════════════════════════════════════════════════════════════════════════
#  DECK 2 — SUPER ADMIN GUIDE (for the founder)
# ═══════════════════════════════════════════════════════════════════════════

def build_admin_guide():
    prs = new_prs()
    DOC_TITLE = "UU Aviation Recruitment 2026  ·  Super Admin Guide"
    slides = []

    # ── Cover ────────────────────────────────────────────────────────────
    slides.append(title_slide(
        prs,
        kicker="SUPER ADMIN GUIDE  ·  EDITION 2026.05",
        title="The full control panel,\nin your hand.",
        subtitle=(
            "Everything you (and only you) need to know to run, monitor, "
            "recover and finalise the drive. Covers: setup, day-of operations, "
            "monitoring, incident response, and post-drive cleanup."
        ),
        audience="Bhupender · Founder · Super Admin",
    ))

    # ── TOC ──────────────────────────────────────────────────────────────
    s = blank_slide(prs)
    page_header(s, "WHAT'S IN THIS GUIDE",
                "Six parts. Pre-drive (read once), day-of (open this deck on a tablet).")
    parts = [
        ("01", "Pre-drive setup",      "T-7 days · staff accounts · rooms · DNS"),
        ("02", "The 5 admin screens",  "What each one shows · what each one is for"),
        ("03", "Day-of operations",    "Watch list · interventions · escalation rules"),
        ("04", "Incident response",    "When the site goes down · when DB pauses · when email fails"),
        ("05", "Security & compliance","Audit log · rate limits · secret rotation"),
        ("06", "Post-drive wrap-up",   "Backup · archive · re-deployment for next drive"),
    ]
    y = Inches(1.45)
    for i, (num, title, hint) in enumerate(parts):
        row_y = y + Inches(0.62) * i
        add_rect(s, Inches(0.55), row_y, Inches(0.55), Inches(0.50), NAVY)
        add_text(s, Inches(0.55), row_y, Inches(0.55), Inches(0.50),
                 num, size=14, bold=True, color=SURFACE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Georgia")
        add_text(s, Inches(1.3), row_y + Inches(0.02),
                 Inches(7.5), Inches(0.3), title,
                 size=14, bold=True, color=NAVY, font="Georgia")
        add_text(s, Inches(1.3), row_y + Inches(0.28),
                 Inches(11), Inches(0.3), hint,
                 size=11, color=MUTED)
    slides.append(s)

    # ─── PART 01 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 1, "Pre-drive setup",
        "T-7 to T-1 days. Everything that has to be true before doors open."))

    slides.append(steps_slide(
        prs, "PART 01  ·  PRE-DRIVE", "T-7 days checklist",
        [
            ("1.  Confirm the domain works",
             "Open https://careers.ews.aero in incognito. Landing page loads in < 1 sec. "
             "If not — Vercel deployment didn't promote or DNS is stale."),
            ("2.  Log in with your password",
             "/admin/login → bhupender@eliteworldservices.com + password. "
             "If your password is forgotten: ADMIN_EMERGENCY_PASSWORD env var (rotate after use)."),
            ("3.  Create one staff account per recruiter/desk operator",
             "/admin/users → Add new staff form → fill email, name, role, set initial password (min 8 chars). "
             "Email them the password OOB (WhatsApp/SMS — never share via the same email)."),
            ("4.  Assign each recruiter to a room",
             "Same /admin/users page → in their row, pick a room from the 'Room' dropdown. "
             "Only RECRUITER rows show this dropdown. Saves instantly."),
            ("5.  Print the QR poster",
             "/admin/qr-poster → Print → A3 colour preferred. Make 5+ copies for entrances."),
            ("6.  Print the runbook",
             "/admin/runbook → Print → A4. One per desk. Same info as the playbook deck."),
            ("7.  Do a full dry-run with 2 fake students",
             "Register two test emails. Walk one through to SELECTED, one through to REJECTED. "
             "Confirm both got the result email."),
            ("8.  Take a fresh DB backup",
             "Run scripts/backup.sh locally. Saves a pg_dump to ~/uu-backups/. Pre-drive snapshot."),
        ]))

    slides.append(three_card_slide(
        prs, "PART 01  ·  ROOM ASSIGNMENT", "How the new dropdown works (added 2026-05-28)",
        [
            {"title": "WHAT IT FIXES", "accent": GREEN, "body": [
                "Before: recruiters got 'No room assigned' and you had to write SQL.",
                "Now: dropdown on /admin/users — pick a room → saved instantly → audit logged.",
                "Auto-revalidates /recruiter + /display so the recruiter's screen flips immediately.",
            ]},
            {"title": "RULES", "accent": BLUE, "body": [
                "One recruiter = one room. Picking a new room clears the old one.",
                "One room = one recruiter. The dropdown hides rooms already taken by someone else.",
                "Only RECRUITER role gets the dropdown. Other roles see '—'.",
                "Picking '— no room —' unassigns them (e.g. recruiter went home sick).",
            ]},
            {"title": "CURRENT SEED", "accent": NAVY, "body": [
                "R1 · Boardroom",
                "R2 · Cabin Crew",
                "R3 · Ground Staff",
                "R4 · Customer Service",
                "Need more rooms? See Part 04 — add via DB or seed script.",
            ]},
        ]))

    # ─── PART 02 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 2, "The 5 admin screens",
        "Memorise where each lives. You will switch between them all day."))

    slides.append(three_card_slide(
        prs, "PART 02  ·  YOUR 5 SCREENS", "Operational dashboards",
        [
            {"title": "/admin   (Overview)", "accent": NAVY, "body": [
                "5 tiles: Registered · Waiting · In progress · Completed · Selected.",
                "Auto-refresh on every page load.",
                "Watch the gap between Waiting and Completed — it's your throughput pulse.",
                "Loads in ~350 ms (single GROUP BY query, no N+1).",
            ]},
            {"title": "/admin/queue   (Live queue)", "accent": GREEN, "body": [
                "Full waiting list across all rooms — operator's home screen.",
                "Filter chips: All · Waiting · Called · Skipped · Done · No-show.",
                "Per-row: Recall · No-show · Search by token.",
                "Use this as your war-room view.",
            ]},
            {"title": "/admin/students   (Roster)", "accent": BLUE, "body": [
                "Search by name/email/token. Filter by status/course/semester.",
                "Click any row → full profile with resume + photo + interview history.",
                "'Export CSV' top right — full roster, all filters applied.",
            ]},
        ]))

    slides.append(three_card_slide(
        prs, "PART 02  ·  YOUR 5 SCREENS (cont.)", "Outreach + governance",
        [
            {"title": "/admin/emails", "accent": GREEN, "body": [
                "Bulk announcement composer.",
                "Pick audience (status × course × semester) → Preview → Type count (if >25) → Send.",
                "Sends via Resend SMTP. Free tier = 3000/month, 100/day.",
                "Use AFTER the drive for selection/rejection batches.",
            ]},
            {"title": "/admin/users", "accent": AMBER, "body": [
                "CRUD for staff accounts.",
                "Add / reset password / activate / deactivate.",
                "Room dropdown for recruiters.",
                "You CANNOT delete yourself or deactivate yourself.",
            ]},
            {"title": "/admin/analytics", "accent": NAVY, "body": [
                "Tiles + SVG charts: status distribution, hourly throughput, rooms breakdown.",
                "Pure-server SVG — no JS chart library. Loads instantly.",
                "Open at the end of each hour for a pulse check.",
            ]},
        ],
        footer_text="Also: /admin/runbook (printable cheat sheet) · /admin/qr-poster (printable QR + venue info)"))

    # ─── PART 03 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 3, "Day-of operations",
        "Hour-by-hour. What you should be doing at each phase of the drive."))

    slides.append(steps_slide(
        prs, "PART 03  ·  YOUR DAY", "A 6-hour drive in 6 phases",
        [
            ("T-60 min  ·  Pre-flight",
             "Run the 8-step pre-flight (Part 01 of the team playbook). "
             "Confirm every staff member can log in. Test one fake registration."),
            ("T-0 to T+30 min  ·  Arrival rush",
             "Watch /admin tiles. Registered should climb fast. "
             "If it stalls → check the registration page is reachable from a student phone (test on 4G, not your wifi)."),
            ("T+30 min to T+5 hr  ·  Steady state",
             "Open 3 tabs: /admin, /admin/queue, /display (kiosk view). "
             "Refresh every 10 min. Eyeball anomalies."),
            ("T+5 hr  ·  Tail end",
             "Waiting count starts dropping. Walk the floor — any students who haven't been called? "
             "Help desk operator clear no-shows."),
            ("T+6 hr  ·  Doors close",
             "Confirm 0 WAITING tokens. Email Manager queues selection blast. Take final screenshot."),
            ("T+6 hr +30 min  ·  Hand-over",
             "Run backup script. Archive Supabase storage bucket. Disable staff accounts you won't reuse."),
        ]))

    slides.append(two_col_slide(
        prs, "PART 03  ·  ESCALATION", "When to act — when to wait",
        left_title="ACT IMMEDIATELY",
        left_items=[
            "Registered count went BACKWARDS — DB integrity bug, stop the drive",
            "Site returns 5xx for > 60 sec — promote previous deployment on Vercel",
            "Supabase dashboard shows 'paused' — un-pause from Supabase project page",
            "Bulk email tool shows 'over quota' — Resend free tier exhausted, contact me",
            "A recruiter says 'I can see the wrong student' — refresh; if persists, log them out and back in",
        ],
        right_title="WAIT & MONITOR",
        right_items=[
            "Single student's submit failed once — they'll retry, not your problem",
            "Admin page took 4 sec to load on first nav — cold start, normal",
            "Display board didn't update for 5 sec — that's the poll interval, normal",
            "One photo failed to upload — student can submit without it (fixable after)",
            "Audit log empty for last hour — only writes on actions; quiet hour = nothing to log",
        ],
        footer_text="Rule of thumb: 1 anomaly = monitor · 3 same anomalies = escalate · 1 systemic anomaly = intervene NOW"))

    # ─── PART 04 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 4, "Incident response",
        "Three classes of incident. Six playbooks. Drilled, tested, ready."))

    slides.append(issue_fix_slide(
        prs, "PART 04  ·  INFRASTRUCTURE INCIDENTS",
        "Vercel / Supabase / Resend",
        [
            ("Site returns 500 on every URL",
             "Bad deployment promoted",
             "Vercel → Deployments → last green → 'Promote to production'. ~30 sec to recover."),
            ("Site loads but every page shows 'DB error'",
             "Supabase paused (free tier) OR connection limit hit",
             "Supabase → Project → Resume. If limit hit: wait 60 sec — our pool cap recycles."),
            ("Confirmation emails not arriving",
             "Resend quota burned OR sender domain DNS broke",
             "Resend dashboard → Logs (delivery / bounces). If quota: switch off bulk send, transactional still works."),
            ("/display TV froze on one tile",
             "Browser tab crashed locally",
             "Refresh the TV. Server data is fine — this is client-only."),
        ]))

    slides.append(issue_fix_slide(
        prs, "PART 04  ·  DATA + ACCESS INCIDENTS",
        "Login locks, data leaks, password issues",
        [
            ("You're locked out of your own admin",
             "10 failed login attempts triggered the 15-min ban",
             "Wait 15 min. OR: log in via emergency password (set in Vercel env). Always rotate it after use."),
            ("Recruiter says 'someone else's data appeared'",
             "Stale tab or proxy cache",
             "Have them log out, clear cookies for careers.ews.aero, log back in. If still: audit log → check actorId on the affected row."),
            ("A student says 'I got someone else's admit card'",
             "Should be impossible — each card now has a per-student token in the URL",
             "Verify by opening the URL — should require ?t=TOKEN. If not: a stale email link was clicked. Resend their email from /admin/students."),
            ("Suspicion: someone scraped the student list",
             "Check audit log for unusual frequency on /api/admin/students.csv",
             "All admin routes now require a valid session — proxy enforces this at the edge. Rotate session secret + force all logins out (delete all sessions in DB)."),
        ]))

    # ─── PART 05 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 5, "Security & compliance",
        "What's protecting the data — and what you should do to keep it that way."))

    slides.append(three_card_slide(
        prs, "PART 05  ·  WHAT'S LIVE", "Defence layers in production",
        [
            {"title": "AT THE EDGE", "accent": NAVY, "body": [
                "HSTS (2-year) — browsers refuse HTTP",
                "X-Frame-Options DENY — can't embed in iframe",
                "X-Content-Type-Options nosniff",
                "Referrer-Policy strict-origin",
                "Permissions-Policy — cam/mic/geo OFF",
                "proxy.ts guards /admin/* + /recruiter + /api/admin/*",
            ]},
            {"title": "AT THE APP", "accent": GREEN, "body": [
                "bcrypt-hashed staff passwords (cost 10)",
                "HMAC-signed session cookies (12 h)",
                "IP rate limit: 10 logins / 5 min / 15 min ban",
                "Email rate limit: 8 attempts / 5 min / account",
                "Failed logins audited with IP + reason",
                "Recruiter actions room-ownership checked",
            ]},
            {"title": "AT THE DATA", "accent": BLUE, "body": [
                "Per-student admit-card token (no enumeration)",
                "Success page gated on token (or admin session)",
                "CSV export injection-safe (= + - @ prefix → quote)",
                "Bulk-email typed confirmation > 25 recipients",
                "Audit log on every staff write action",
                "studentId never trusted from client",
            ]},
        ]))

    slides.append(steps_slide(
        prs, "PART 05  ·  YOUR HOUSEKEEPING", "Five rules — non-negotiable",
        [
            ("1.  Never share a staff password over the same email",
             "Use WhatsApp/SMS for initial passwords. Staff resets their own at first login (post-MVP)."),
            ("2.  Rotate ADMIN_EMERGENCY_PASSWORD after any use",
             "Generate new value (32+ chars) → Vercel dashboard → env vars → save → redeploy."),
            ("3.  Run scripts/backup.sh after every drive",
             "pg_dump + storage bucket archive. Keep 30 days."),
            ("4.  Disable accounts after the drive ends",
             "/admin/users → Deactivate each external recruiter. Reactivate next drive — history preserved."),
            ("5.  Review the audit log monthly",
             "DB: SELECT action, COUNT(*) FROM \"AuditLog\" WHERE \"createdAt\" > NOW() - INTERVAL '30 days' GROUP BY action."),
        ]))

    # ─── PART 06 ─────────────────────────────────────────────────────────
    slides.append(section_slide(
        prs, 6, "Post-drive wrap-up",
        "From last student walking out to next drive's prep — clean and archived."))

    slides.append(steps_slide(
        prs, "PART 06  ·  POST-DRIVE", "8-step wrap-up",
        [
            ("1.  Confirm 0 WAITING + 0 IN_PROGRESS",
             "Anything stuck → no-show it. The system is now quiet."),
            ("2.  Run selection / rejection bulk emails",
             "/admin/emails. Send selection first (good news). Rejection second, with a kind close."),
            ("3.  Export the master CSV",
             "/admin/students → Export. Save as drive-YYYY-MM-DD.csv in your Drive folder."),
            ("4.  Run scripts/backup.sh",
             "pg_dump + S3 sync of Supabase storage bucket. ~2 min for 700 students."),
            ("5.  Take final analytics screenshot",
             "/admin/analytics → full-page screenshot for the post-drive report."),
            ("6.  Deactivate one-off recruiter accounts",
             "/admin/users → Deactivate. Keep your super admin + perma staff active."),
            ("7.  Archive the Vercel deployment URL",
             "Save the URL of today's deployment in case you need to roll back to it later."),
            ("8.  Pause Supabase project (if no traffic expected)",
             "Saves the free-tier hours. Resume 1 day before next drive."),
        ]))

    slides.append(two_col_slide(
        prs, "PART 06  ·  PREPARING THE NEXT DRIVE", "Reuse the same system, reset the data",
        left_title="DON'T REBUILD — RESET",
        left_items=[
            "Wipe Student / Token / InterviewLog tables (NOT User / Room)",
            "Keep staff accounts active — passwords still valid",
            "Update room display names if course mix changes",
            "Update brand colours in src/app/globals.css if needed",
        ],
        right_title="UPDATE BEFORE GO-LIVE",
        right_items=[
            "Reprint QR poster with new event date (auto-detects from runbook)",
            "Re-do dry-run with 2 fake students end-to-end",
            "Bump APP_URL env if domain changes",
            "Take a fresh DB backup BEFORE wiping (the prior drive's data)",
        ],
        footer_text="The system is built once, run many times. Average prep time for drive #2 = 1 hour."))

    # ── Closing ─────────────────────────────────────────────────────────
    slides.append(closing_slide(
        prs,
        "You built this. You own this.",
        [
            "Everything in this guide is one click away on careers.ews.aero.",
            "",
            "If a screen confuses you mid-drive: open /admin/runbook on a phone — it's the one-page version.",
            "If the runbook doesn't help: this deck is your full reference.",
            "If this deck doesn't help: roll back the deployment on Vercel and call the dev. The system stays up.",
        ],
        contacts=[
            "Live URL: careers.ews.aero",
            "Vercel: dashboard.vercel.com",
            "Supabase: supabase.com/dashboard",
            "Resend: resend.com/emails",
        ]))

    # ── Footers ──────────────────────────────────────────────────────────
    total = len(slides)
    for i, s in enumerate(slides[2:], start=2):
        footer(s, i + 1, total, DOC_TITLE)

    out = OUT_DIR / "UU-Aviation-Super-Admin-Guide.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    a = build_playbook()
    b = build_admin_guide()
    print(f"Wrote: {a}")
    print(f"Wrote: {b}")
